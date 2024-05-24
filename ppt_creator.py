from pptx import Presentation

# pip install python-pptx

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Air-speed Velocity of Unladen Swallows"

tf = body_shape.text_frame

for i in range(5):
    p = tf.add_paragraph()

    p.text = f"bullet {i}"
    p.level = 0


prs.save("test.pptx")
